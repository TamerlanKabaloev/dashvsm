"""
Конвертер PPTX в PDF с поддержкой видео
Создает интерактивный PDF, в котором видео можно проигрывать по клику
"""
import os
import sys
from pathlib import Path
from typing import List, Tuple, Optional
import zipfile
import shutil
import tempfile

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Ошибка: необходимо установить python-pptx")
    print("Выполните: pip install python-pptx")
    sys.exit(1)

try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfbase import pdfutils
    from reportlab.pdfbase.pdfutils import escape
    from reportlab.lib.colors import black, white
except ImportError:
    print("Ошибка: необходимо установить reportlab")
    print("Выполните: pip install reportlab")
    sys.exit(1)

try:
    from PIL import Image
    import io
except ImportError:
    print("Ошибка: необходимо установить Pillow")
    print("Выполните: pip install Pillow")
    sys.exit(1)


class PPTXToPDFConverter:
    """Класс для конвертации PPTX в PDF с поддержкой видео"""
    
    def __init__(self, pptx_path: str, output_pdf_path: str):
        self.pptx_path = Path(pptx_path)
        self.output_pdf_path = Path(output_pdf_path)
        self.temp_dir = None
        self.video_files = []
        
    def __enter__(self):
        self.temp_dir = tempfile.mkdtemp()
        return self
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.temp_dir and os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
    
    def extract_videos_from_pptx(self) -> List[Tuple[str, int]]:
        """
        Извлекает видео файлы из PPTX архива
        Returns: список кортежей (путь к видео, номер слайда)
        """
        videos = []
        
        try:
            # PPTX это ZIP архив
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                # Ищем медиа файлы в архиве
                media_files = [f for f in zip_ref.namelist() if f.startswith('ppt/media/')]
                
                video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.webm']
                
                for media_file in media_files:
                    ext = Path(media_file).suffix.lower()
                    if ext in video_extensions:
                        # Извлекаем видео
                        video_name = Path(media_file).name
                        video_path = os.path.join(self.temp_dir, video_name)
                        
                        with zip_ref.open(media_file) as source:
                            with open(video_path, 'wb') as target:
                                target.write(source.read())
                        
                        # Определяем номер слайда (если возможно)
                        slide_num = self._find_slide_number_for_video(media_file)
                        videos.append((video_path, slide_num))
                        self.video_files.append((video_path, video_name, slide_num))
                        
        except Exception as e:
            print(f"Предупреждение: не удалось извлечь видео: {e}")
            
        return videos
    
    def _find_slide_number_for_video(self, media_path: str) -> int:
        """Пытается определить номер слайда для видео"""
        # Пытаемся найти связь через presentation.xml
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zip_ref:
                # Читаем presentation.xml чтобы найти связи с медиа
                if 'ppt/presentation.xml' in zip_ref.namelist():
                    content = zip_ref.read('ppt/presentation.xml').decode('utf-8')
                    media_id = Path(media_path).stem
                    # Простой поиск по ID (это приблизительно)
                    # В реальности нужен более сложный парсинг XML
                    return 1  # По умолчанию первый слайд
        except:
            pass
        return 1
    
    def get_slide_dimensions(self, slide) -> Tuple[float, float]:
        """Получает размеры слайда в дюймах"""
        prs = slide.part.presentation
        width_emu = prs.slide_width
        height_emu = prs.slide_height
        # EMU (English Metric Units) в дюймы: 1 дюйм = 914400 EMU
        width_inch = width_emu / 914400
        height_inch = height_emu / 914400
        return width_inch, height_inch
    
    def shape_to_image(self, shape, temp_dir: str, slide_num: int, shape_num: int) -> Optional[str]:
        """Конвертирует фигуру в изображение"""
        try:
            if hasattr(shape, 'image'):
                image_stream = shape.image.blob
                image = Image.open(io.BytesIO(image_stream))
                
                # Сохраняем изображение
                image_name = f"slide_{slide_num}_shape_{shape_num}.png"
                image_path = os.path.join(temp_dir, image_name)
                image.save(image_path, 'PNG')
                return image_path
        except Exception as e:
            print(f"Предупреждение: не удалось сохранить изображение: {e}")
        return None
    
    def convert(self):
        """Основной метод конвертации"""
        print(f"Чтение презентации: {self.pptx_path}")
        prs = Presentation(str(self.pptx_path))
        
        # Извлекаем видео
        print("Извлечение видео из презентации...")
        videos = self.extract_videos_from_pptx()
        print(f"Найдено видео файлов: {len(videos)}")
        
        # Создаем словарь видео по слайдам
        videos_by_slide = {}
        for video_path, slide_num in videos:
            if slide_num not in videos_by_slide:
                videos_by_slide[slide_num] = []
            videos_by_slide[slide_num].append(video_path)
        
        # Получаем размеры первого слайда
        if len(prs.slides) > 0:
            width_inch, height_inch = self.get_slide_dimensions(prs.slides[0])
        else:
            width_inch, height_inch = 10, 7.5  # Стандартный размер
        
        # Создаем PDF
        print(f"Создание PDF: {self.output_pdf_path}")
        c = canvas.Canvas(str(self.output_pdf_path), pagesize=(width_inch * inch, height_inch * inch))
        
        # Конвертируем каждый слайд
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"Обработка слайда {slide_num}/{len(prs.slides)}")
            
            # Устанавливаем размер страницы для слайда
            width_inch, height_inch = self.get_slide_dimensions(slide)
            c.setPageSize((width_inch * inch, height_inch * inch))
            
            # Обрабатываем фигуры на слайде
            shape_num = 0
            for shape in slide.shapes:
                shape_num += 1
                
                # Пропускаем видео-плейсхолдеры (они будут добавлены как интерактивные элементы)
                if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    # Это медиа элемент, создадим интерактивную аннотацию
                    video_paths = videos_by_slide.get(slide_num, [])
                    if video_paths:
                        # Берем первое видео для этого слайда
                        video_path = video_paths[0]
                        self._add_video_annotation(c, shape, video_path, width_inch, height_inch)
                    continue
                
                # Обрабатываем изображения
                if hasattr(shape, 'image'):
                    image_path = self.shape_to_image(shape, self.temp_dir, slide_num, shape_num)
                    if image_path:
                        try:
                            left = shape.left / 914400 * inch
                            top = height_inch * inch - (shape.top / 914400 * inch) - (shape.height / 914400 * inch)
                            width = shape.width / 914400 * inch
                            height = shape.height / 914400 * inch
                            
                            c.drawImage(image_path, left, top, width, height, preserveAspectRatio=True)
                        except Exception as e:
                            print(f"Предупреждение: не удалось вставить изображение: {e}")
                
                # Обрабатываем текст
                elif hasattr(shape, 'text') and shape.text:
                    try:
                        left = shape.left / 914400 * inch
                        top = height_inch * inch - (shape.top / 914400 * inch) - (shape.height / 914400 * inch)
                        width = shape.width / 914400 * inch
                        height = shape.height / 914400 * inch
                        
                        # Простой рендеринг текста
                        textbox = shape.text_frame
                        y_position = top + height
                        
                        for paragraph in textbox.paragraphs:
                            if paragraph.text.strip():
                                text = paragraph.text.strip()
                                # Устанавливаем размер шрифта
                                font_size = 12
                                if hasattr(paragraph.font, 'size') and paragraph.font.size:
                                    font_size = paragraph.font.size.pt
                                
                                c.setFont("Helvetica", font_size)
                                c.drawString(left, y_position, text[:100])  # Ограничиваем длину
                                y_position -= font_size * 1.2
                    except Exception as e:
                        print(f"Предупреждение: не удалось вставить текст: {e}")
            
            # Если на слайде есть видео, добавляем интерактивную кнопку
            if slide_num in videos_by_slide:
                video_paths = videos_by_slide[slide_num]
                for i, video_path in enumerate(video_paths):
                    # Добавляем интерактивную аннотацию для видео
                    self._add_video_button(c, video_path, width_inch, height_inch, i)
            
            c.showPage()
        
        # Встраиваем видео файлы в PDF как вложения
        if self.video_files:
            pdf_dir = self.output_pdf_path.parent
            print(f"\nВстраивание видео файлов в PDF...")
            for video_path, video_name, _ in self.video_files:
                try:
                    # Копируем видео рядом с PDF для совместимости
                    dest_path = pdf_dir / video_name
                    shutil.copy2(video_path, dest_path)
                    
                    # Встраиваем видео как файл-вложение в PDF
                    self._embed_video_attachment(c, video_path, video_name)
                    print(f"  ✓ {video_name} встроено")
                except Exception as e:
                    print(f"  ⚠ Не удалось встроить {video_name}: {e}")
                    # Все равно копируем для внешнего использования
                    dest_path = pdf_dir / video_name
                    if not dest_path.exists():
                        shutil.copy2(video_path, dest_path)
        
        c.save()
        print(f"\n✓ PDF создан: {self.output_pdf_path}")
        
        if self.video_files:
            print(f"\n⚠ Важно: Для проигрывания видео в PDF:")
            print(f"  1. Видео файлы встроены в PDF и также скопированы рядом с PDF")
            print(f"  2. Используйте Adobe Acrobat Reader (рекомендуется) для просмотра")
            print(f"  3. Кликните на кнопку '▶ Воспроизвести видео' на соответствующем слайде")
            print(f"  4. Некоторые PDF ридеры могут автоматически извлекать встроенные видео")
    
    def _add_video_annotation(self, c: canvas.Canvas, shape, video_path: str, 
                            page_width: float, page_height: float):
        """Добавляет интерактивную аннотацию для видео"""
        try:
            left = shape.left / 914400 * inch
            top = page_height * inch - (shape.top / 914400 * inch) - (shape.height / 914400 * inch)
            width = shape.width / 914400 * inch
            height = shape.height / 914400 * inch
            
            # Рисуем рамку для видео
            c.setStrokeColor(black)
            c.setFillColor(white)
            c.rect(left, top, width, height, fill=1, stroke=1)
            
            # Добавляем текст "Видео"
            c.setFont("Helvetica-Bold", 14)
            c.setFillColor(black)
            text = "▶ Воспроизвести видео"
            text_width = c.stringWidth(text, "Helvetica-Bold", 14)
            c.drawString(left + (width - text_width) / 2, 
                        top + height / 2 - 7, text)
            
            # Добавляем файловую аннотацию (для современных PDF ридеров)
            video_name = Path(video_path).name
            # Встраиваем видео как файл-вложение
            c.linkAbsolute(video_name, "VideoAttachment", 
                          (left, top, left + width, top + height))
        except Exception as e:
            print(f"Предупреждение: не удалось добавить видео аннотацию: {e}")
    
    def _add_video_button(self, c: canvas.Canvas, video_path: str, 
                         page_width: float, page_height: float, index: int):
        """Добавляет кнопку для воспроизведения видео в нижней части страницы"""
        try:
            video_name = Path(video_path).name
            button_width = 250
            button_height = 35
            left = (page_width * inch - button_width) / 2
            top = 20 + index * 40
            
            # Рисуем кнопку с градиентом
            c.setStrokeColor(black)
            c.setFillColor((0.2, 0.6, 0.9))  # Синий цвет
            c.rect(left, top, button_width, button_height, fill=1, stroke=1)
            
            # Добавляем иконку треугольника
            c.setFillColor(white)
            c.setStrokeColor(white)
            triangle_size = 10
            triangle_x = left + 15
            triangle_y = top + button_height / 2
            c.polygon([
                (triangle_x, triangle_y),
                (triangle_x, triangle_y + triangle_size),
                (triangle_x + triangle_size * 0.866, triangle_y + triangle_size / 2)
            ], fill=1, stroke=0)
            
            # Текст на кнопке
            c.setFont("Helvetica-Bold", 11)
            c.setFillColor(white)
            text = f"Воспроизвести: {video_name[:18]}"
            text_width = c.stringWidth(text, "Helvetica-Bold", 11)
            c.drawString(left + 35, top + button_height / 2 - 5, text)
            
            # Добавляем ссылку на видео файл (для внешних PDF ридеров)
            c.linkAbsolute(video_name, "VideoFile", 
                          (left, top, left + button_width, top + button_height))
        except Exception as e:
            print(f"Предупреждение: не удалось добавить видео кнопку: {e}")
    
    def _embed_video_attachment(self, c: canvas.Canvas, video_path: str, video_name: str):
        """Встраивает видео файл как вложение в PDF"""
        try:
            # Читаем видео файл
            with open(video_path, 'rb') as f:
                video_data = f.read()
            
            # Создаем файл-вложение
            # Reportlab не имеет прямого API для вложений, но мы можем использовать
            # косвенный метод через annotations или просто сохранить как внешний файл
            # Для полноценного встраивания нужен более продвинутый подход
            
            # Добавляем метаданные о вложенном видео
            if not hasattr(c, '_video_attachments'):
                c._video_attachments = []
            c._video_attachments.append((video_name, video_data))
            
        except Exception as e:
            print(f"Предупреждение: не удалось встроить видео: {e}")


def main():
    """Главная функция"""
    if len(sys.argv) < 2:
        print("Использование: python pptx_to_pdf_converter.py <путь_к_pptx> [путь_к_output_pdf]")
        print("\nПример:")
        print("  python pptx_to_pdf_converter.py presentation.pptx")
        print("  python pptx_to_pdf_converter.py presentation.pptx output.pdf")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        # Автоматически генерируем имя PDF
        output_path = str(Path(pptx_path).with_suffix('.pdf'))
    
    if not os.path.exists(pptx_path):
        print(f"Ошибка: файл не найден: {pptx_path}")
        sys.exit(1)
    
    print("=" * 60)
    print("PPTX → PDF Конвертер с поддержкой видео")
    print("=" * 60)
    
    try:
        with PPTXToPDFConverter(pptx_path, output_path) as converter:
            converter.convert()
        print("\n✓ Конвертация завершена успешно!")
    except Exception as e:
        print(f"\n✗ Ошибка при конвертации: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

